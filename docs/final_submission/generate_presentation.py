from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

OUTPUT = "CSUF_Scheduler_Final_Presentation.pptx"

prs = Presentation()

# Theme-like constants
BLUE = RGBColor(26, 72, 126)
DARK = RGBColor(32, 32, 32)
ACCENT = RGBColor(0, 141, 213)


def add_title_slide(title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    _style_title(slide.shapes.title)


def add_bullets_slide(title: str, bullets: list[str], subtitle: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    _style_title(slide.shapes.title)

    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    if subtitle:
        p = body.paragraphs[0]
        p.text = subtitle
        p.font.bold = True
        p.font.size = Pt(22)
        p.font.color.rgb = ACCENT

    for item in bullets:
        p = body.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = DARK


def _style_title(title_shape) -> None:
    tf = title_shape.text_frame
    for p in tf.paragraphs:
        for run in p.runs:
            run.font.color.rgb = BLUE
            run.font.bold = True
            run.font.size = Pt(36)


# 1. Title
add_title_slide(
    "CSUF Scheduler",
    "Master's Project Defense\nDevelopment Project\nSpring 2026",
)

# 2. Problem and Motivation
add_bullets_slide(
    "Problem and Motivation",
    [
        "Manual scheduling creates conflicts, inequity, and high admin overhead.",
        "Student-staff operations require transparent, auditable shift management.",
        "Need for web + mobile access with secure role-based controls.",
    ],
)

# 3. Objectives
add_bullets_slide(
    "Project Objectives",
    [
        "Build production-focused scheduling platform for CSUF operations.",
        "Implement role-based workflows: admin, supervisor, student, FTO.",
        "Provide AI-assisted schedule generation with manual oversight.",
        "Deliver measurable quality via testing, linting, and security checks.",
    ],
)

# 4. Architecture
add_bullets_slide(
    "System Architecture",
    [
        "Backend: Flask + SQLAlchemy + Alembic migration pipeline.",
        "Service layer isolates scheduling logic (OR-Tools engine).",
        "API layer exposes users, shifts, scheduler, reports, conflicts, audit.",
        "Frontend strategy: Flask templates + Vue modernization + React Native mobile MVP.",
    ],
)

# 5. Data and API Design
add_bullets_slide(
    "Data and API Design",
    [
        "Normalized domain entities: users, shifts, availability, timesheets, audit.",
        "Session-based auth with CSRF protection for mutating requests.",
        "Health endpoint for operational checks: /api/health.",
        "Report exports support JSON and CSV consumption.",
    ],
)

# 6. Security and Reliability
add_bullets_slide(
    "Security and Reliability",
    [
        "Session hardening, secure headers, CSRF enforcement, rate limiting.",
        "Pre-commit quality gates: Black, isort, Flake8, Bandit.",
        "Dependency audit process documented with mitigation tracking.",
        "Operational runbooks: deployment, monitoring, incident response, smoke tests.",
    ],
)

# 7. Implemented Features
add_bullets_slide(
    "Implemented Features",
    [
        "Shift lifecycle: create, assign, auto-assign, swap, delete.",
        "Availability and notifications workflows.",
        "Scheduler endpoints: generate, inspect, apply, cancel schedules.",
        "Reports and conflicts endpoints for weekly planning visibility.",
    ],
)

# 8. Mobile MVP (Phase 7)
add_bullets_slide(
    "Mobile MVP (Phase 7)",
    [
        "React Native + Expo app integrated with existing backend APIs.",
        "Screens: Login, Schedule, Swap, Availability, Notifications, Settings.",
        "Preserved backend auth contract (session + CSRF) for consistency.",
        "Validated startup and diagnostics with Expo tooling.",
    ],
)

# 9. Validation Results
add_bullets_slide(
    "Validation Results",
    [
        "Health endpoint: healthy, database OK.",
        "Smoke flow: login, fetch users/shifts, create shift, verify created shift.",
        "Automated tests: 64/64 passed.",
        "Expo doctor: 16/16 checks passed; Metro startup validated.",
    ],
)

# 10. Deployment Readiness
add_bullets_slide(
    "Deployment Readiness",
    [
        "Docker compose production baseline prepared.",
        "Environment templates for staging and production included.",
        "CI/CD workflow configured for quality + deployment pipeline.",
        "Documented post-deployment smoke validation checklist.",
    ],
)

# 11. Challenges and Decisions
add_bullets_slide(
    "Key Challenges and Decisions",
    [
        "Balanced rapid delivery with production hardening requirements.",
        "Chose session auth over JWT to match current architecture and reduce risk.",
        "Kept frontend modernization incremental to preserve compatibility.",
        "Contained optimization complexity within service layer abstractions.",
    ],
)

# 12. Future Work
add_bullets_slide(
    "Future Enhancements",
    [
        "Push notifications and deeper mobile-native capabilities.",
        "Expanded accessibility and usability testing.",
        "Advanced observability with tracing and performance dashboards.",
        "Enterprise controls: MFA/SSO and stronger compliance automation.",
    ],
)

# 13. Conclusion
add_bullets_slide(
    "Conclusion",
    [
        "CSUF Scheduler achieved development-project objectives.",
        "System is validated, documented, and operationally ready for controlled deployment.",
        "Architecture supports continued scaling and feature expansion.",
    ],
)

# 14. Q&A
add_bullets_slide(
    "Questions",
    [
        "Thank you.",
        "I welcome questions on architecture, validation, or deployment strategy.",
    ],
)

prs.save(OUTPUT)
print(f"Created {OUTPUT}")
